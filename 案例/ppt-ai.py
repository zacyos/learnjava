# Import necessary libraries
import openai
import requests
import json
import re
import pptx
from pptx.util import Inches

# Set up OpenAI API credentials
openai.api_key = "YOUR_API_KEY"

# Define function to generate text using GPT-3
def generate_text(prompt):
    response = openai.Completion.create(
        engine="davinci",
        prompt=prompt,
        temperature=0.5,
        max_tokens=1024,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    text = response.choices[0].text
    return text

# Define function to create optimized PowerPoint presentation
def create_ppt(prompt):
    # Generate text using GPT-3
    text = generate_text(prompt)
    
    # Clean up text
    text = re.sub(r'\n+', '\n', text)
    text = text.strip()
    
    # Split text into slides
    slides = text.split('\n\n')
    
    # Create PowerPoint presentation
    prs = pptx.Presentation()
    
    # Add slides to presentation
    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, body = slide_text.split('\n', 1)
        slide.shapes.title.text = title
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = body
        tf.fit_text(font_family='Calibri', max_size=18)
        tf.auto_size = True
    
    # Save presentation
    prs.save('optimized_presentation.pptx')
    
    return 'optimized_presentation.pptx'

# Call create_ppt function with GPT-3 prompt
create_ppt('Create a PowerPoint presentation on the benefits of using GPT-3 for generating optimized presentations.')